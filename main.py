from pptx import Presentation
from pptx.util import Inches, Pt, Cm
import pptx.enum.text
import pandas as pd
import pptx.enum.shapes
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import ColorFormat, RGBColor
import numpy as np

def addSlide(name, prs):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    addBackground(prs, slide)
    addHeart(slide)
    addName(slide, name)

def addName(slide_name, name):
    left = Inches(2.05)
    top = Inches(3.55)
    width = Cm(15)
    height = Cm(1)
    txBox = slide_name.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = name
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def addHeart(slide):
    shapes = slide.shapes
    left = Inches(0.09)
    top = Inches(0)
    width = Inches(9.82)
    height = Inches(7.5)
    
    shape = shapes.add_shape(MSO_SHAPE.HEART, left, top, width, height)
    fill = shape.fill
    fill.solid()
    
    fill.fore_color.rgb = RGBColor(255,255,255)
   
    line = shape.line
    line.color.rgb = RGBColor(0,0,0)





def addBackground(prs, slide):
    slides = prs.slides
    slide_num = slides.index(slide)+1



    switcher = {
        0: RGBColor(180,167,214),
        1: RGBColor(234,209,220),
        2: RGBColor(255,229,153),
        3: RGBColor(207,226,243),
        4: RGBColor(234,153,153),
        5: RGBColor(217,210,233),
        6: RGBColor(244,204,204),
        7: RGBColor(159,197,232),
        8: RGBColor(255,242,204),
        9: RGBColor(213,166,189)
    }


    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = switcher.get(slide_num % 10)



def make():
    prs = Presentation()
    pathXL  ="/Users/sunilkumar/Downloads/Enrolled Students SP21_1.xlsx"
    df= pd.read_excel(pathXL, engine = 'openpyxl')
    
    
    for index,row in df.iterrows():
        first_name = row['Campus First Name']
        last_name = row['Last Name']

       
        if str(first_name) != 'nan' and str(last_name) != 'nan':
            name = f'{first_name} {last_name}'
            addSlide(name, prs)
    
    prs.save('test.pptx')

make()